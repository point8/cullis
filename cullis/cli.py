"""Main entrypoint and CLI for cullis
"""

import os
import click
import shutil
import logging
import tempfile

from pptx import Presentation
from pptx.util import Inches
try:
    from wand.image import Image
except:
    raise Exception('MAGICK_HOME not set, see README for instructions')
from pdfrw import PdfReader, PdfWriter

from cullis.logger import configure_logger
from cullis.version import __version__


NAME = 'cullis'
BLANK_SLIDE = 6
EMU_PER_INCH = 914400
LOCAL_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_DIR = os.path.join(LOCAL_DIR, 'templates')

logger = configure_logger()


@click.group(help='Vom Regen in die Traufe')
@click.option('-v', '--verbose', is_flag=True)
@click.version_option(version=__version__, prog_name='cullis')
def cli(verbose):
    if verbose:
        file_handler.setLevel(logging.DEBUG)


@cli.command(help='')
@click.argument('source', required=True, type=click.Path(exists=True))
@click.option('-v', '--verbose', is_flag=True, help='Show debug output')
@click.option('-o', '--out', help='Name of output file (w/o file extension)')
def convert(source, verbose, out):
    if verbose:
        logger.setLevel(logging.DEBUG)

    src = click.format_filename(source)
    out = out or 'converted'

    # let everything happen in tmp dir
    with tempfile.TemporaryDirectory() as tmpdir:

        if os.path.isfile(src) and src.lower().endswith('.pdf'):
            logger.info(f'Reading {src}')
            with open(src, 'rb') as infile:
                pdf = PdfReader(infile)

            logger.info(f'Converting PDF pages to images')
            for idx, page in enumerate(pdf.pages):
                logger.debug(f'Page {idx}')
                # Every page as single PDF file
                writer = PdfWriter()
                writer.addpage(page)
                writer.write(os.path.join(tmpdir, f'conv.{idx:>03}.pdf'))

                # Converting PDF into JPG
                with Image(filename=os.path.join(tmpdir, f'conv.{idx:>03}.pdf')) as img:
                     img.save(filename=os.path.join(tmpdir, f'conv.{idx:>03}.jpg'))
    
                os.remove(os.path.join(tmpdir, f'conv.{idx:>03}.pdf'))

        if os.path.isdir(src):
            logger.info(f'Reading {src} and collect all image files')
            for root, dirs, files in os.walk(src):
                for f in files:
                    # copy all image files to tmp dir
                    if f.lower().endswith(('.png', '.jpg', '.jpeg')):
                        shutil.copy(os.path.join(root, f), tmpdir)
    
        # Now start the conversion to pptx
        pptx = Presentation(os.path.join(TEMPLATE_DIR, 'template.pptx'))
        h, w = pptx.slide_height, pptx.slide_width
    
        left = top = Inches(0)
        height = Inches(h / EMU_PER_INCH)
    
        blank_slide_layout = pptx.slide_layouts[BLANK_SLIDE]
    
        # walk through tmp dir and add all image files to presentation
        for root, dirs, files in os.walk(tmpdir):
            logger.info(f'Adding {len(files)} image files to presentation')
            for idx, f in enumerate(sorted(files)):
                if not f.lower().endswith(('.png', '.jpg', '.jpeg')):
                    logger.debug(f'Skipping {f}')
                    continue

                logger.debug(f'Adding {f} file to presentation')
                slide = None
                if idx == 0:
                    slide = pptx.slides[0]
                else:
                    slide = pptx.slides.add_slide(blank_slide_layout)

                slide.shapes.add_picture(os.path.join(tmpdir, f), left, top, height=height)

        out_file = f'{out}.pptx'
        pptx.save(out_file)
        logger.info(f'Saved presentation to {out_file}')


def main():
    cli()
