import os

from pptx import Presentation
from pptx.util import Inches
from wand.image import Image
from pdfrw import PdfReader, PdfWriter


def main():
    INPATH = '/Users/ccauet/Nextcloud/Point 8/Vorträge/2017-11 bmcExchange/2017-11-16_bmcExchange_p8.pdf'

    with open(INPATH, 'rb') as infile:
        pdf = PdfReader(infile)

    print(pdf.pages[0])

    writer = PdfWriter()
    writer.addpage(pdf.pages[0])
    writer.write('tmp.pdf')

    pptx = Presentation()

    from wand.image import Image
    # Converting PDF into JPG
    with Image(filename="tmp.pdf") as img:
         img.save(filename="temp.jpg")

    pptx = Presentation()
    blank_slide_layout = pptx.slide_layouts[6]
    slide = pptx.slides.add_slide(blank_slide_layout)

    left = top = Inches(1)
    # pic = slide.shapes.add_picture('temp.jpg', left, top)

    left = Inches(5)
    height = Inches(5.5)
    pic = slide.shapes.add_picture('temp.jpg', left, top, height=height)

    pptx.save('test.pptx')


if __name__ == '__main__':
    main()
